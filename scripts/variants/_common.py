"""Shared helpers + constants for the alternative poster variants.

The Polished Classic variant in ``scripts/build_poster.py`` is the source of
truth for poster *content* (TITLE, ABSTRACT, METHODS_BULLETS, RESULTS_TABLE,
TAKEHOME_TOP/MID/BOT, CONCLUSION_BULLETS, SUPPORTING_BULLETS, REFERENCES).
The variant scripts under this package import those constants verbatim so
every variant tells the same story; only the *visual presentation* differs.

This module re-exports those constants, adds a few helpers that every
variant needs (rectangle, text-box, font-cascaded paragraph, deterministic
zip normalisation, LibreOffice preview rendering), and exposes the fixed
author line that every variant must use on its byline.
"""

from __future__ import annotations

import datetime as _dt
import os
import shutil
import subprocess
import sys
import tempfile
import zipfile

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
FIG_DIR = os.path.join(REPO_ROOT, "figures")
POSTER_DIR = os.path.join(REPO_ROOT, "poster")
VARIANTS_DIR = os.path.join(POSTER_DIR, "variants")
SCRIPTS_DIR = os.path.join(REPO_ROOT, "scripts")

# Pull the canonical content from the Polished Classic build script.
if SCRIPTS_DIR not in sys.path:
    sys.path.insert(0, SCRIPTS_DIR)

from build_poster import (  # noqa: E402
    TITLE,
    HEADER_KICKER,
    ABSTRACT,
    QUESTION,
    BACKGROUND_BULLETS,
    METHODS_BULLETS,
    RESULTS_TABLE,
    TAKEHOME_TOP,
    TAKEHOME_MID,
    TAKEHOME_BOT,
    CONCLUSION_BULLETS,
    SUPPORTING_BULLETS,
    REFERENCES,
)

# Every variant uses the same byline string (orchestrator requirement).
VARIANT_AUTHOR_LINE = (
    "Ivan Del Rio \u00b7 Anish Kondamadugula  |  COGS 109 (Mukamel) \u00b7 "
    "Spring 2026  |  UC San Diego"
)

# 48" x 36" landscape, identical to the Polished Classic.
SLIDE_W_IN = 48.0
SLIDE_H_IN = 36.0

# Figure paths (all variants share these).
FIG_HEADLINE = os.path.join(FIG_DIR, "11_knn_k_sweep.png")
FIG_BACKGROUND = os.path.join(FIG_DIR, "08_label_autocorrelation.png")
FIG_SUPPORTING = os.path.join(FIG_DIR, "14_cv_comparison_three_way.png")
FIG_CONFUSION = os.path.join(FIG_DIR, "15_confusion_matrices.png")


def add_rect(slide, x, y, w, h, fill, line=None, line_width_pt=0.75):
    """Add a filled rectangle, optionally with a thin keyline."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width_pt)
    shape.shadow.inherit = False
    return shape


def add_line_rect(slide, x, y, w, h, fill):
    """Solid coloured rectangle used as a horizontal / vertical rule."""
    return add_rect(slide, x, y, w, h, fill)


def add_text(slide, x, y, w, h, *, fill=None, anchor=MSO_ANCHOR.TOP, margin=0.1):
    """Add a text box with sensible defaults and return ``(box, frame)``."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin / 2)
    tf.margin_bottom = Inches(margin / 2)
    tf.paragraphs[0].text = ""
    return tb, tf


def set_para(p, text, *, size, bold=False, italic=False, color, align=PP_ALIGN.LEFT,
             font="Calibri", spacing=None):
    """Set the first run of ``p`` to a styled single line of text."""
    p.alignment = align
    if spacing is not None:
        p.space_after = Pt(spacing)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color


def add_paragraph(tf, text, *, size, bold=False, italic=False, color,
                  align=PP_ALIGN.LEFT, font="Calibri", spacing=6, first=False):
    """Append a paragraph to ``tf`` (or fill the first one) with styled text."""
    if first or (not tf.paragraphs[0].runs and tf.paragraphs[0].text == ""):
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    set_para(p, text, size=size, bold=bold, italic=italic, color=color,
             align=align, font=font, spacing=spacing)
    return p


def bullets(tf, items, *, size, color, font="Calibri", glyph="\u2022  ",
            leading=4, bold=False):
    """Render ``items`` as bulleted paragraphs inside ``tf``."""
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(leading)
        run = p.add_run()
        run.text = glyph + item
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color


def insert_image_fit(slide, path, x, y, max_w, max_h):
    """Insert ``path`` into ``slide`` centred inside the (max_w, max_h) box.

    The image is scaled with aspect ratio preserved.  Returns the actual
    (x, y, w, h) box the image occupies so the caller can position a caption.
    """
    from PIL import Image as _Image
    with _Image.open(path) as im:
        iw, ih = im.size
    target_w = max_w
    target_h = int(target_w * ih / iw)
    if target_h > max_h:
        target_h = max_h
        target_w = int(target_h * iw / ih)
    off_x = x + (max_w - target_w) // 2
    off_y = y + (max_h - target_h) // 2
    slide.shapes.add_picture(path, off_x, off_y, width=target_w, height=target_h)
    return off_x, off_y, target_w, target_h


def normalize_pptx_zip(path):
    """Re-pack ``path`` with deterministic ZIP order + fixed timestamps.

    Mirrors ``scripts/build_poster.py::_normalize_pptx_zip`` so variant
    builds are also byte-stable across runs.
    """
    fixed_date = (2026, 5, 28, 0, 0, 0)
    with zipfile.ZipFile(path, "r") as src:
        infos = sorted(src.infolist(), key=lambda zi: zi.filename)
        payloads = [(zi.filename, src.read(zi.filename)) for zi in infos]
    tmp_path = path + ".tmp"
    with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as dst:
        for name, data in payloads:
            zi = zipfile.ZipInfo(filename=name, date_time=fixed_date)
            zi.compress_type = zipfile.ZIP_DEFLATED
            zi.external_attr = 0o644 << 16
            dst.writestr(zi, data)
    os.replace(tmp_path, path)


def set_core_properties(prs, *, subject):
    """Pin core properties so two builds produce a byte-identical archive."""
    cp = prs.core_properties
    cp.author = "BariBariGood"
    cp.title = TITLE
    cp.subject = subject
    fixed = _dt.datetime(2026, 5, 28, 0, 0, 0)
    cp.created = fixed
    cp.modified = fixed
    cp.last_modified_by = "BariBariGood"


def render_preview_png(pptx_path, png_path):
    """Render a PNG preview next to ``pptx_path`` using LibreOffice headless.

    If ``soffice`` is not on PATH or the conversion fails, write a small
    placeholder PNG so the build doesn't crash; the .pptx is the canonical
    deliverable either way.
    """
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice is not None:
        with tempfile.TemporaryDirectory(prefix="variant_render_") as td:
            try:
                subprocess.run(
                    [soffice, "--headless", "--convert-to", "png",
                     "--outdir", td, pptx_path],
                    check=True, capture_output=True, timeout=180,
                )
                base = os.path.splitext(os.path.basename(pptx_path))[0]
                rendered = os.path.join(td, base + ".png")
                if os.path.exists(rendered):
                    shutil.copyfile(rendered, png_path)
                    return png_path
            except (subprocess.CalledProcessError, subprocess.TimeoutExpired):
                pass
    _placeholder_png(png_path)
    return png_path


def _placeholder_png(png_path):
    """Write a tiny PIL placeholder when LibreOffice is unavailable."""
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (1440, 1080), (240, 240, 240))
    d = ImageDraw.Draw(img)
    d.text((40, 40),
           "Preview unavailable (LibreOffice not found). Open the .pptx "
           "directly.", fill=(20, 20, 20))
    img.save(png_path)


__all__ = [
    "TITLE", "HEADER_KICKER", "ABSTRACT", "QUESTION",
    "BACKGROUND_BULLETS", "METHODS_BULLETS", "RESULTS_TABLE",
    "TAKEHOME_TOP", "TAKEHOME_MID", "TAKEHOME_BOT",
    "CONCLUSION_BULLETS", "SUPPORTING_BULLETS", "REFERENCES",
    "VARIANT_AUTHOR_LINE",
    "SLIDE_W_IN", "SLIDE_H_IN",
    "FIG_HEADLINE", "FIG_BACKGROUND", "FIG_SUPPORTING", "FIG_CONFUSION",
    "REPO_ROOT", "FIG_DIR", "POSTER_DIR", "VARIANTS_DIR",
    "add_rect", "add_line_rect", "add_text", "set_para", "add_paragraph",
    "bullets", "insert_image_fit",
    "normalize_pptx_zip", "set_core_properties", "render_preview_png",
    "RGBColor", "MSO_ANCHOR", "PP_ALIGN", "Inches", "Pt",
]
