#!/usr/bin/env python3
"""Variant 4 — Minimal / Whitespace (Apple Keynote / Stripe / Linear spirit).

A 48" x 36" landscape poster built around aggressive negative space.
Pure white background, near-black text, a single electric-blue accent,
light gray separators.  One modern sans throughout, aggressive type
hierarchy: a 120pt title, 24pt all-caps tracked-out gray section labels,
26pt body, and a 200pt one-line take-home stat.  Layout is a vertical
narrative (not a 3-column grid) that reads top-to-bottom like a Stripe
blog post: title \u2192 question \u2192 hero chart \u2192 single-line take-home
\u2192 two-column methods / conclusions \u2192 thin bottom strip with
references and tiny supporting figures.
"""

from __future__ import annotations

import os

from pptx import Presentation

from scripts.variants._common import (
    ABSTRACT, BACKGROUND_BULLETS, CONCLUSION_BULLETS, FIG_BACKGROUND,
    FIG_HEADLINE, FIG_SUPPORTING, Inches, MSO_ANCHOR, PP_ALIGN, Pt,
    QUESTION, REFERENCES, RGBColor, RESULTS_TABLE, METHODS_BULLETS,
    SLIDE_H_IN, SLIDE_W_IN, SUPPORTING_BULLETS, TAKEHOME_BOT, TAKEHOME_MID,
    TAKEHOME_TOP, TITLE, VARIANTS_DIR, VARIANT_AUTHOR_LINE, add_paragraph,
    add_rect, add_text, bullets, insert_image_fit, normalize_pptx_zip,
    render_preview_png, set_core_properties, set_para,
)

# ----- palette -----
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x0A, 0x0A, 0x0A)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GRAY = RGBColor(0xE5, 0xE7, 0xEB)
MUTED = RGBColor(0x6B, 0x72, 0x80)

SANS = "Arial"
SERIF = "Georgia"


def _hairline(slide, x, y, w, color=GRAY, weight=0.03):
    add_rect(slide, x, y, w, Inches(weight), color)


def _section_label(slide, x, y, w, h, label):
    """All-caps tracked-out gray section label."""
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    tracked = " ".join(list(label.upper()))
    set_para(tf.paragraphs[0], tracked, size=22, bold=True, color=MUTED,
             align=PP_ALIGN.LEFT, font=SANS)


def _bullets(slide, x, y, w, h, items, *, size=22, leading=10):
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, margin=0.0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(leading)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.name = SANS
        run.font.color.rgb = INK


def _trim(text, max_words):
    """Trim a string to roughly ``max_words`` words for whitespace variant."""
    parts = text.split()
    if len(parts) <= max_words:
        return text
    return " ".join(parts[:max_words]) + " \u2026"


def _hero_takehome(slide, x, y, w, h):
    """Three-line take-home statement, biggest type on the poster.

    The brief asks for a single 200pt line, but the actual content (
    "78% honest. 97% leaky. The difference is leakage.") is too wide to
    fit at that size on a 48" canvas with generous side margins.  We
    keep the same words and tone but stack the three clauses vertically
    so the typography stays Apple/Stripe oversized.
    """
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    p0 = tf.paragraphs[0]
    set_para(p0, "78% honest.", size=160, bold=True, color=BLUE,
             align=PP_ALIGN.CENTER, font=SANS)
    p1 = tf.add_paragraph()
    p1.space_before = Pt(4)
    set_para(p1, "97% leaky.", size=160, bold=True, color=MUTED,
             align=PP_ALIGN.CENTER, font=SANS)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    set_para(p2, "The difference is leakage.", size=80, bold=True, color=INK,
             align=PP_ALIGN.CENTER, font=SANS)


def build():
    os.makedirs(VARIANTS_DIR, exist_ok=True)
    out_pptx = os.path.join(VARIANTS_DIR, "variant_4_minimal_whitespace.pptx")
    out_png = os.path.join(VARIANTS_DIR, "variant_4_minimal_whitespace_preview.png")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Pure white background.
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE)

    # Generous side margins for the Stripe / Apple aesthetic.
    side = Inches(3.0)
    inner_w = prs.slide_width - 2 * side

    # ----- TITLE BAND -----
    # Tiny tracked-out kicker
    tb, tf = add_text(slide, side, Inches(0.9), inner_w, Inches(0.5),
                      anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    set_para(tf.paragraphs[0],
             " ".join(list("COGS 109  \u00b7  SPRING 2026")),
             size=18, bold=True, color=MUTED,
             align=PP_ALIGN.LEFT, font=SANS)

    # Huge left-aligned title (auto-wraps to 2 lines on a long line)
    tb, tf = add_text(slide, side, Inches(1.5), inner_w, Inches(5.4),
                      anchor=MSO_ANCHOR.TOP, margin=0.0)
    set_para(tf.paragraphs[0], TITLE, size=100, bold=True, color=INK,
             align=PP_ALIGN.LEFT, font=SANS)

    # Thin blue rule under the title
    _hairline(slide, side, Inches(7.1), Inches(10.0), color=BLUE, weight=0.06)

    # Byline below rule
    tb, tf = add_text(slide, side, Inches(7.35), inner_w, Inches(0.55),
                      anchor=MSO_ANCHOR.TOP, margin=0.0)
    set_para(tf.paragraphs[0], VARIANT_AUTHOR_LINE, size=22, bold=False,
             color=MUTED, align=PP_ALIGN.LEFT, font=SANS)

    # ----- QUESTION (one-sentence serif italic, centred, lots of whitespace) -----
    q_text = (
        "How does the choice of cross-validation scheme change the "
        "result of an otherwise-identical KNN model-selection procedure?"
    )
    tb, tf = add_text(slide, side + Inches(2.0), Inches(8.5),
                      inner_w - Inches(4.0), Inches(2.1),
                      anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    set_para(tf.paragraphs[0], q_text, size=36, italic=True, color=INK,
             align=PP_ALIGN.CENTER, font=SERIF)

    # ----- HEADLINE FIGURE (centered with generous side margins) -----
    fig_top = Inches(11.0)
    fig_max_h = Inches(10.6)
    fig_box_w = inner_w  # already has ~6" left+right whitespace via `side`
    insert_image_fit(slide, FIG_HEADLINE, side, fig_top, fig_box_w, fig_max_h)
    # Single-line caption beneath, muted gray
    tb, tf = add_text(slide, side, Inches(22.0), inner_w, Inches(0.6),
                      anchor=MSO_ANCHOR.TOP, margin=0.0)
    set_para(tf.paragraphs[0],
             "KNN model selection under three CV schemes \u2014 "
             "all three pick k = 1.",
             size=20, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER, font=SANS)

    # ----- HERO TAKE-HOME -----
    _hero_takehome(slide, side, Inches(22.9), inner_w, Inches(6.4))

    # ----- TWO-COLUMN METHODS / CONCLUSIONS -----
    col2_top = Inches(29.6)
    col2_h = Inches(3.6)
    col_w = (inner_w - Inches(1.5)) // 2
    col1_x = side
    col2_x = side + col_w + Inches(1.5)

    # Vertical blue accent line between columns (thin)
    add_rect(slide, col2_x - Inches(0.85), col2_top + Inches(0.2),
             Inches(0.04), col2_h - Inches(0.4), BLUE)

    _section_label(slide, col1_x, col2_top, col_w, Inches(0.4), "Methods")
    methods_trimmed = [
        "KNN with Euclidean distance in z-scored 14-D channel space.",
        "Sweep k over log-spaced {1\u2026201}; pick k maximising mean 5-fold CV accuracy.",
        "Three CV schemes: shuffled (leaky), naive blocked, stratified blocked (honest).",
    ]
    _bullets(slide, col1_x, col2_top + Inches(0.55), col_w, col2_h - Inches(0.55),
             methods_trimmed, size=20, leading=8)

    _section_label(slide, col2_x, col2_top, col_w, Inches(0.4), "Conclusions")
    conclusions_trimmed = [
        "All three schemes pick k = 1, but reported accuracy spans 47 pp.",
        "Honest accuracy: 77.8% \u00b1 2.7% \u2014 well above the 55.12% baseline.",
        "+19.5 pp gap is scheme-choice leakage, not signal. KNN is uniquely vulnerable.",
    ]
    _bullets(slide, col2_x, col2_top + Inches(0.55), col_w, col2_h - Inches(0.55),
             conclusions_trimmed, size=20, leading=8)

    # ----- THIN BOTTOM STRIP: references + tiny supporting figs -----
    _hairline(slide, side, Inches(33.4), inner_w, color=GRAY, weight=0.03)

    strip_top = Inches(33.65)
    # Tiny insets on the left (autocorr) and right (supporting)
    inset_w = Inches(3.0)
    inset_h = Inches(1.7)
    if os.path.exists(FIG_BACKGROUND):
        insert_image_fit(slide, FIG_BACKGROUND, side, strip_top, inset_w, inset_h)
    if os.path.exists(FIG_SUPPORTING):
        insert_image_fit(slide, FIG_SUPPORTING,
                         side + inner_w - inset_w, strip_top,
                         inset_w, inset_h)
    # Tiny caption strip under each inset
    tb, tf = add_text(slide, side, strip_top + inset_h + Inches(0.0),
                      inset_w, Inches(0.4), anchor=MSO_ANCHOR.TOP, margin=0.05)
    set_para(tf.paragraphs[0],
             "Autocorrelation across lags",
             size=12, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER, font=SANS)
    tb, tf = add_text(slide, side + inner_w - inset_w,
                      strip_top + inset_h + Inches(0.0),
                      inset_w, Inches(0.4), anchor=MSO_ANCHOR.TOP, margin=0.05)
    set_para(tf.paragraphs[0],
             "Three-way CV comparison \u2014 LDA / KNN / PCA\u2192LDA / PCR",
             size=12, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER, font=SANS)

    # References live in the central gap between the insets
    refs_x = side + inset_w + Inches(0.8)
    refs_w = inner_w - 2 * (inset_w + Inches(0.8))
    _section_label(slide, refs_x, strip_top, refs_w, Inches(0.35),
                   "References")
    tb, tf = add_text(slide, refs_x, strip_top + Inches(0.45),
                      refs_w, Inches(2.0), anchor=MSO_ANCHOR.TOP, margin=0.0)
    for i, ref in enumerate(REFERENCES):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = ref
        run.font.size = Pt(11)
        run.font.name = SANS
        run.font.color.rgb = MUTED

    set_core_properties(prs,
                        subject="COGS 109 EEG poster \u2014 Minimal Whitespace variant")
    prs.save(out_pptx)
    normalize_pptx_zip(out_pptx)
    render_preview_png(out_pptx, out_png)
    return out_pptx, out_png


if __name__ == "__main__":
    p, png = build()
    print(f"wrote {p}")
    print(f"wrote {png}")
