#!/usr/bin/env python3
"""Variant 2 — Editorial Hero (NYT / 538 / Pudding.cool data-viz spirit).

A 48" x 36" landscape poster where the headline numbers ARE the visual.
Cream off-white background, dark charcoal text, a single burnt-orange
accent, and muted slate-gray for context.  Top third is a stat strip
(97% leaky / 78% honest / 50% punitive); middle third is the KNN k-sweep
figure full-width with thin orange rules above and below; bottom third
is a clean 3-column layout for methods, conclusions, and references
with the autocorrelation and supporting-evidence figures as small
inset images in the side columns.
"""

from __future__ import annotations

import os

from pptx import Presentation

from scripts.variants._common import (
    ABSTRACT, BACKGROUND_BULLETS, CONCLUSION_BULLETS, FIG_BACKGROUND,
    FIG_HEADLINE, FIG_SUPPORTING, Inches, MSO_ANCHOR, PP_ALIGN, Pt,
    QUESTION, REFERENCES, RGBColor, RESULTS_TABLE, METHODS_BULLETS,
    SLIDE_H_IN, SLIDE_W_IN, SUPPORTING_BULLETS, TAKEHOME_BOT, TAKEHOME_MID,
    TAKEHOME_TOP, TITLE, VARIANTS_DIR, VARIANT_AUTHOR_LINE, add_paragraph,
    add_rect, add_text, bullets, insert_image_fit, normalize_pptx_zip,
    render_preview_png, set_core_properties, set_para,
)

# ----- palette -----
BG_CREAM = RGBColor(0xFA, 0xF7, 0xF2)
CHARCOAL = RGBColor(0x1A, 0x1A, 0x1A)
ORANGE = RGBColor(0xD9, 0x77, 0x57)
SLATE = RGBColor(0x5B, 0x67, 0x70)
LIGHT_SLATE = RGBColor(0x9A, 0xA2, 0xA8)

# Font families with conservative fallbacks (LibreOffice substitutes as needed).
SANS_BOLD = "Arial"   # paired with bold=True
SERIF = "Georgia"     # body serif
SERIF_ITALIC = "Georgia"  # italic flag is set per-run


def _kicker(slide, x, y, w, h, text):
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE, margin=0.05)
    set_para(tf.paragraphs[0], text.upper(), size=22, bold=True, color=SLATE,
             align=PP_ALIGN.LEFT, font=SANS_BOLD)


def _title(slide, x, y, w, h):
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE, margin=0.05)
    set_para(tf.paragraphs[0], TITLE, size=84, bold=True, color=CHARCOAL,
             align=PP_ALIGN.LEFT, font=SANS_BOLD)


def _byline(slide, x, y, w, h):
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE, margin=0.05)
    set_para(tf.paragraphs[0], VARIANT_AUTHOR_LINE, size=24, bold=False,
             color=SLATE, align=PP_ALIGN.LEFT, font=SERIF, italic=True)


def _orange_rule(slide, x, y, w, *, weight=0.07):
    add_rect(slide, x, y, w, Inches(weight), ORANGE)


def _stat_column(slide, x, y, w, h, numeral, label, hero=False):
    """One column of the stat strip: giant numeral + italic serif label."""
    num_h = int(h * 0.74)
    label_h = h - num_h
    num_box, num_tf = add_text(slide, x, y, w, num_h,
                               anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    num_color = ORANGE if hero else LIGHT_SLATE
    # Hero numeral is noticeably bigger; muted ones recede.
    num_size = 440 if hero else 320
    set_para(num_tf.paragraphs[0], numeral, size=num_size, bold=True,
             color=num_color, align=PP_ALIGN.CENTER, font=SANS_BOLD)
    lab_box, lab_tf = add_text(slide, x, y + num_h, w, label_h,
                               anchor=MSO_ANCHOR.TOP, margin=0.0)
    set_para(lab_tf.paragraphs[0], label, size=44, italic=True, bold=False,
             color=CHARCOAL if hero else SLATE,
             align=PP_ALIGN.CENTER, font=SERIF)


def _pull_quote(slide, x, y, w, h):
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE, margin=0.2)
    quote = "\u201c{} {} {}\u201d".format(TAKEHOME_TOP, TAKEHOME_MID, TAKEHOME_BOT)
    set_para(tf.paragraphs[0], quote, size=46, italic=True, bold=False,
             color=CHARCOAL, align=PP_ALIGN.CENTER, font=SERIF)


def _headline_figure(slide, x, y, w, h):
    insert_image_fit(slide, FIG_HEADLINE, x, y, w, h)


def _figure_caption(slide, x, y, w, text, *, size=18, color=None):
    tb, tf = add_text(slide, x, y, w, Inches(0.55),
                      anchor=MSO_ANCHOR.TOP, margin=0.05)
    set_para(tf.paragraphs[0], text, size=size, italic=True, bold=False,
             color=color or SLATE, align=PP_ALIGN.CENTER, font=SERIF)


def _section_label(slide, x, y, w, label):
    """Small caps tracked-out section label in orange."""
    tb, tf = add_text(slide, x, y, w, Inches(0.5),
                      anchor=MSO_ANCHOR.TOP, margin=0.0)
    set_para(tf.paragraphs[0], label.upper(), size=22, bold=True, color=ORANGE,
             align=PP_ALIGN.LEFT, font=SANS_BOLD)
    # Thin horizontal rule beneath the label
    add_rect(slide, x, y + Inches(0.55), w, Inches(0.025), ORANGE)


def _bullet_column(slide, x, y, w, h, items, *, size=18):
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, margin=0.0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "\u2014  " + item   # em-dash bullet, editorial flavour
        run.font.size = Pt(size)
        run.font.name = SERIF
        run.font.color.rgb = CHARCOAL


def _ref_column(slide, x, y, w, h, refs, *, size=14):
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, margin=0.0)
    for i, ref in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = ref
        run.font.size = Pt(size)
        run.font.name = SERIF
        run.font.color.rgb = SLATE


def _by_the_numbers_strip(slide, x, y, w):
    """One-line italic micro-caption summarising the headline numbers."""
    parts = [
        f"shuffled (leaky) {RESULTS_TABLE[1][2]}",
        f"naive blocked {RESULTS_TABLE[2][2]}",
        f"stratified blocked (honest) {RESULTS_TABLE[3][2]}",
    ]
    line = "  \u2022  ".join(parts) + "  (all schemes pick k = 1)"
    tb, tf = add_text(slide, x, y, w, Inches(0.55),
                      anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    set_para(tf.paragraphs[0], line, size=20, italic=True, bold=False,
             color=SLATE, align=PP_ALIGN.CENTER, font=SERIF)


def _abstract_lede(slide, x, y, w, h):
    """Editorial lede paragraph below the pull quote — uses verbatim ABSTRACT."""
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, margin=0.2)
    set_para(tf.paragraphs[0], ABSTRACT, size=22, italic=False, bold=False,
             color=CHARCOAL, align=PP_ALIGN.LEFT, font=SERIF)


def _research_question(slide, x, y, w, h):
    """Italic strapline above the stat strip, framed as 'THE QUESTION'."""
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE, margin=0.2)
    p0 = tf.paragraphs[0]
    set_para(p0, "THE QUESTION    ", size=20, bold=True, color=ORANGE,
             align=PP_ALIGN.CENTER, font=SANS_BOLD)
    run = p0.add_run()
    run.text = QUESTION
    run.font.size = Pt(22)
    run.font.italic = True
    run.font.name = SERIF
    run.font.color.rgb = CHARCOAL


def build():
    os.makedirs(VARIANTS_DIR, exist_ok=True)
    out_pptx = os.path.join(VARIANTS_DIR, "variant_2_editorial_hero.pptx")
    out_png = os.path.join(VARIANTS_DIR, "variant_2_editorial_hero_preview.png")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Cream background
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, BG_CREAM)

    # ----- HEADER (top edge) -----
    side_margin = Inches(2.0)
    inner_w = prs.slide_width - 2 * side_margin

    _kicker(slide, side_margin, Inches(0.7), inner_w, Inches(0.55),
            "UC San Diego  /  Department of Cognitive Science")
    _title(slide, side_margin, Inches(1.4), inner_w, Inches(2.4))
    _byline(slide, side_margin, Inches(3.9), inner_w, Inches(0.55))
    _orange_rule(slide, side_margin, Inches(4.7), inner_w, weight=0.06)

    # ----- THE QUESTION strap-line -----
    _research_question(slide, side_margin, Inches(4.9), inner_w, Inches(0.85))

    # ----- STAT STRIP -----
    strip_top = Inches(5.9)
    strip_h = Inches(6.8)
    gap = Inches(0.6)
    stat_w = (inner_w - 2 * gap) // 3
    cols_x = [
        side_margin,
        side_margin + stat_w + gap,
        side_margin + 2 * (stat_w + gap),
    ]
    _stat_column(slide, cols_x[0], strip_top, stat_w, int(strip_h),
                 "97%", "shuffled / leaky", hero=False)
    _stat_column(slide, cols_x[1], strip_top, stat_w, int(strip_h),
                 "78%", "stratified blocked / honest", hero=True)
    _stat_column(slide, cols_x[2], strip_top, stat_w, int(strip_h),
                 "50%", "naive blocked / punitive", hero=False)

    _by_the_numbers_strip(slide, side_margin, strip_top + strip_h + Inches(0.15),
                          inner_w)

    # ----- PULL QUOTE -----
    pq_top = Inches(13.6)
    pq_h = Inches(2.2)
    _pull_quote(slide, side_margin + Inches(1.5),
                pq_top, inner_w - Inches(3.0), pq_h)

    # ----- HEADLINE FIGURE -----
    _orange_rule(slide, side_margin, Inches(16.1), inner_w, weight=0.05)
    fig_top = Inches(16.4)
    fig_h = Inches(13.4)
    _headline_figure(slide, side_margin + Inches(0.5), fig_top,
                     inner_w - Inches(1.0), fig_h)
    # Editorial side annotations: small italic notes flanking the figure.
    side_note_w = Inches(8.0)
    # Left annotation
    tb_l, tf_l = add_text(slide, side_margin + Inches(0.2),
                          fig_top + Inches(1.0), side_note_w, Inches(3.4),
                          anchor=MSO_ANCHOR.TOP, margin=0.1)
    set_para(tf_l.paragraphs[0],
             "Shuffled 5-fold CV", size=20, bold=True, color=ORANGE,
             align=PP_ALIGN.LEFT, font=SANS_BOLD)
    p_l = tf_l.add_paragraph()
    p_l.space_before = Pt(4)
    set_para(p_l,
             "Sits at 0.97 across every k. Looks dazzling, but the random "
             "shuffle leaks adjacent samples (lag-1 autocorr \u2248 0.997) "
             "into both train and test.",
             size=15, italic=True, color=CHARCOAL,
             align=PP_ALIGN.LEFT, font=SERIF)
    # Right annotation (hero — the honest line)
    right_x = side_margin + inner_w - side_note_w - Inches(0.2)
    tb_r, tf_r = add_text(slide, right_x, fig_top + Inches(1.0),
                          side_note_w, Inches(3.4),
                          anchor=MSO_ANCHOR.TOP, margin=0.1)
    set_para(tf_r.paragraphs[0],
             "Stratified blocked CV", size=20, bold=True, color=ORANGE,
             align=PP_ALIGN.LEFT, font=SANS_BOLD)
    p_r = tf_r.add_paragraph()
    p_r.space_before = Pt(4)
    set_para(p_r,
             "0.778 \u00b1 0.027 at k=1 — honest, time-aware, well above "
             "the 0.5512 majority-class floor. This is the number to "
             "quote.",
             size=15, italic=True, color=CHARCOAL,
             align=PP_ALIGN.LEFT, font=SERIF)
    # Bottom-left annotation for the naive blocked line
    tb_bl, tf_bl = add_text(slide, side_margin + Inches(0.2),
                            fig_top + Inches(8.4), side_note_w, Inches(3.0),
                            anchor=MSO_ANCHOR.TOP, margin=0.1)
    set_para(tf_bl.paragraphs[0],
             "Naive blocked CV", size=20, bold=True, color=ORANGE,
             align=PP_ALIGN.LEFT, font=SANS_BOLD)
    p_bl = tf_bl.add_paragraph()
    p_bl.space_before = Pt(4)
    set_para(p_bl,
             "0.500 \u00b1 0.107: too punitive — 5 contiguous chunks let "
             "class proportion drift across folds, so the mean drops far "
             "below baseline.",
             size=15, italic=True, color=CHARCOAL,
             align=PP_ALIGN.LEFT, font=SERIF)
    # Bottom-right pull-out
    tb_br, tf_br = add_text(slide, right_x,
                            fig_top + Inches(8.4), side_note_w, Inches(3.0),
                            anchor=MSO_ANCHOR.TOP, margin=0.1)
    set_para(tf_br.paragraphs[0],
             "The headline gap", size=20, bold=True, color=ORANGE,
             align=PP_ALIGN.LEFT, font=SANS_BOLD)
    p_br = tf_br.add_paragraph()
    p_br.space_before = Pt(4)
    set_para(p_br,
             "0.973 \u2212 0.778 = +19.5 percentage points. That gap is "
             "leakage attributable to scheme choice on an autocorrelated "
             "single-subject recording \u2014 not a better model.",
             size=15, italic=True, color=CHARCOAL,
             align=PP_ALIGN.LEFT, font=SERIF)

    cap_y = fig_top + fig_h + Inches(0.1)
    _figure_caption(slide, side_margin, cap_y, inner_w,
                    "Fig. 1.  KNN mean accuracy \u00b1 std-dev vs k (log scale) "
                    "under three CV schemes \u2014 all three pick k = 1, but "
                    "the picked accuracies span 47 percentage points.")
    _orange_rule(slide, side_margin, Inches(30.45), inner_w, weight=0.05)

    # ----- BOTTOM 3-COL -----
    bot_top = Inches(30.75)
    col_w = (inner_w - 2 * gap) // 3
    bot_cols = [
        side_margin,
        side_margin + col_w + gap,
        side_margin + 2 * (col_w + gap),
    ]

    # --- Column 1: METHODS + autocorr inset ---
    _section_label(slide, bot_cols[0], bot_top, col_w, "Methods")
    _bullet_column(slide, bot_cols[0], bot_top + Inches(0.75),
                   col_w, Inches(2.6), METHODS_BULLETS, size=13)
    inset_y = bot_top + Inches(3.4)
    inset_h = Inches(2.0)
    if os.path.exists(FIG_BACKGROUND):
        insert_image_fit(slide, FIG_BACKGROUND,
                         bot_cols[0], inset_y, col_w, inset_h)
    _figure_caption(slide, bot_cols[0], inset_y + inset_h + Inches(0.05), col_w,
                    "Fig. 2.  Label autocorrelation across lags \u2014 "
                    "motivates blocked CV.", size=11)

    # --- Column 2: CONCLUSIONS + supporting inset ---
    _section_label(slide, bot_cols[1], bot_top, col_w, "Conclusions")
    _bullet_column(slide, bot_cols[1], bot_top + Inches(0.75),
                   col_w, Inches(2.6), CONCLUSION_BULLETS, size=13)
    if os.path.exists(FIG_SUPPORTING):
        insert_image_fit(slide, FIG_SUPPORTING,
                         bot_cols[1], inset_y, col_w, inset_h)
    _figure_caption(slide, bot_cols[1], inset_y + inset_h + Inches(0.05), col_w,
                    "Fig. 3.  Three-way CV comparison \u2014 KNN has the "
                    "largest leakage gap.", size=11)

    # --- Column 3: SUPPORTING + REFERENCES ---
    _section_label(slide, bot_cols[2], bot_top, col_w, "Supporting evidence")
    _bullet_column(slide, bot_cols[2], bot_top + Inches(0.75),
                   col_w, Inches(2.3), SUPPORTING_BULLETS, size=11)
    _section_label(slide, bot_cols[2], bot_top + Inches(3.25), col_w,
                   "References")
    _ref_column(slide, bot_cols[2], bot_top + Inches(4.0),
                col_w, Inches(1.8), REFERENCES, size=10)

    # ----- FOOTER signature -----
    foot_y = Inches(35.5)
    tb, tf = add_text(slide, side_margin, foot_y, inner_w, Inches(0.4),
                      anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    set_para(tf.paragraphs[0],
             "github.com/BariBariGood/cogs109-eeg-eye-state-classifier",
             size=13, italic=True, bold=False, color=SLATE,
             align=PP_ALIGN.LEFT, font=SERIF)
    tb2, tf2 = add_text(slide, side_margin, foot_y, inner_w, Inches(0.4),
                        anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    set_para(tf2.paragraphs[0],
             "COGS 109 \u2014 Modelling & Data Analysis \u2014 Spring 2026",
             size=13, italic=True, bold=False, color=SLATE,
             align=PP_ALIGN.RIGHT, font=SERIF)

    set_core_properties(prs, subject="COGS 109 EEG poster \u2014 Editorial Hero variant")
    prs.save(out_pptx)
    normalize_pptx_zip(out_pptx)
    render_preview_png(out_pptx, out_png)
    return out_pptx, out_png


if __name__ == "__main__":
    p, png = build()
    print(f"wrote {p}")
    print(f"wrote {png}")
