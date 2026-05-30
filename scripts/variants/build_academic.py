#!/usr/bin/env python3
"""Variant 3 — Academic Classic (NeurIPS / Nature paper-style).

A 48" x 36" landscape poster styled like a calm, deliberate scientific
paper hung on a museum wall.  Parchment cream background, deep ink text,
single brass/ochre accent, small-caps serif section headers underlined
with a thin ochre rule, classical 3-column scientific grid, every panel
edged with a thin ochre keyline, footnoted references in smaller serif
at the bottom.  The methods section reads like a paper's methods —
flowing prose, not bullet points.  Ornamental glyphs (\u00a7) open each
section.
"""

from __future__ import annotations

import os

from pptx import Presentation

from scripts.variants._common import (
    ABSTRACT, BACKGROUND_BULLETS, CONCLUSION_BULLETS, FIG_BACKGROUND,
    FIG_HEADLINE, FIG_SUPPORTING, Inches, MSO_ANCHOR, PP_ALIGN, Pt,
    QUESTION, REFERENCES, RGBColor, RESULTS_TABLE, METHODS_BULLETS,
    SLIDE_H_IN, SLIDE_W_IN, SUPPORTING_BULLETS, TAKEHOME_BOT, TAKEHOME_MID,
    TAKEHOME_TOP, TITLE, VARIANTS_DIR, VARIANT_AUTHOR_LINE, add_paragraph,
    add_rect, add_text, bullets, insert_image_fit, normalize_pptx_zip,
    render_preview_png, set_core_properties, set_para,
)

# ----- palette -----
PARCHMENT = RGBColor(0xF6, 0xEF, 0xE0)
INK = RGBColor(0x20, 0x20, 0x20)
BRASS = RGBColor(0xB5, 0x83, 0x2A)
PANEL = RGBColor(0xED, 0xE5, 0xD0)
MUTED = RGBColor(0x58, 0x4F, 0x3A)

SERIF = "Georgia"

GLYPH = "\u00a7"  # section sign


def _panel(slide, x, y, w, h):
    """Warm-gray panel with a thin brass keyline."""
    add_rect(slide, x, y, w, h, PANEL, line=BRASS, line_width_pt=0.6)


def _section_header(slide, x, y, w, label):
    """Small-caps serif section header with brass underline rule."""
    tb, tf = add_text(slide, x, y, w, Inches(0.7),
                      anchor=MSO_ANCHOR.TOP, margin=0.0)
    p = tf.paragraphs[0]
    # Section glyph in brass, then small-caps title in ink.
    p.alignment = PP_ALIGN.LEFT
    run_g = p.add_run()
    run_g.text = GLYPH + "  "
    run_g.font.size = Pt(28)
    run_g.font.bold = True
    run_g.font.name = SERIF
    run_g.font.color.rgb = BRASS
    run_t = p.add_run()
    # Mimic small-caps by writing uppercase + slight tracking via spaces.
    run_t.text = "  ".join(list(label.upper()))
    run_t.font.size = Pt(24)
    run_t.font.bold = True
    run_t.font.name = SERIF
    run_t.font.color.rgb = INK
    add_rect(slide, x, y + Inches(0.7), w, Inches(0.03), BRASS)


def _body_prose(slide, x, y, w, h, text, *, size=18, justify=True, italic=False):
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, margin=0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.JUSTIFY if justify else PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.italic = italic
    run.font.name = SERIF
    run.font.color.rgb = INK


def _bullets_serif(slide, x, y, w, h, items, *, size=17, leading=4):
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, margin=0.2)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(leading)
        run = p.add_run()
        run.text = "\u25c6  " + item   # ornamental diamond glyph
        run.font.size = Pt(size)
        run.font.name = SERIF
        run.font.color.rgb = INK


def _results_table(slide, x, y, w, h):
    """Render RESULTS_TABLE as a clean serif-styled table."""
    rows, cols = len(RESULTS_TABLE), len(RESULTS_TABLE[0])
    tbl_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl = tbl_shape.table
    col_fracs = [0.46, 0.18, 0.36]
    for i, frac in enumerate(col_fracs):
        tbl.columns[i].width = int(w * frac)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = BRASS
            elif RESULTS_TABLE[r][0].startswith("Stratified blocked"):
                cell.fill.fore_color.rgb = RGBColor(0xE9, 0xD8, 0xA8)
            else:
                cell.fill.fore_color.rgb = PANEL
            tf = cell.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = ""
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = RESULTS_TABLE[r][c]
            run.font.size = Pt(15)
            run.font.bold = (r == 0) or (
                c == 2 and RESULTS_TABLE[r][0].startswith("Stratified blocked")
            )
            run.font.name = SERIF
            run.font.color.rgb = PARCHMENT if r == 0 else INK
    return tbl_shape


def _take_home_block(slide, x, y, w, h):
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE, margin=0.3)
    # Three stacked paragraphs: shuffled / honest / takeaway — biggest emphasis
    # on the middle (honest) line.
    set_para(tf.paragraphs[0], "\u201c" + TAKEHOME_TOP, size=24, italic=True,
             bold=False, color=MUTED, align=PP_ALIGN.CENTER, font=SERIF)
    p1 = tf.add_paragraph()
    p1.space_before = Pt(8)
    set_para(p1, TAKEHOME_MID, size=28, italic=True, bold=True, color=INK,
             align=PP_ALIGN.CENTER, font=SERIF)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    set_para(p2, TAKEHOME_BOT + "\u201d", size=22, italic=True, bold=False,
             color=BRASS, align=PP_ALIGN.CENTER, font=SERIF)


def _figure_with_caption(slide, x, y, w, max_h, path, caption, cap_size=15):
    bx, by, bw, bh = insert_image_fit(slide, path, x, y, w, max_h)
    cap_y = by + bh + Inches(0.05)
    tb, tf = add_text(slide, x, cap_y, w, Inches(0.55),
                      anchor=MSO_ANCHOR.TOP, margin=0.05)
    set_para(tf.paragraphs[0], caption, size=cap_size, italic=True, bold=False,
             color=MUTED, align=PP_ALIGN.CENTER, font=SERIF)
    return bh


def _references_footer(slide, x, y, w, h):
    """Numbered references in small serif at the bottom of the poster."""
    tb, tf = add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, margin=0.1)
    # Section opener as a small-caps mini-header.
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    rg = p0.add_run()
    rg.text = GLYPH + "  R E F E R E N C E S    "
    rg.font.size = Pt(15)
    rg.font.bold = True
    rg.font.name = SERIF
    rg.font.color.rgb = BRASS
    # All references on the same paragraph, separated by em-dashes.
    rt = p0.add_run()
    rt.text = "    " + "    ".join(REFERENCES)
    rt.font.size = Pt(13)
    rt.font.name = SERIF
    rt.font.color.rgb = INK


def build():
    os.makedirs(VARIANTS_DIR, exist_ok=True)
    out_pptx = os.path.join(VARIANTS_DIR, "variant_3_academic_classic.pptx")
    out_png = os.path.join(VARIANTS_DIR, "variant_3_academic_classic_preview.png")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Parchment background
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, PARCHMENT)

    side_margin = Inches(1.4)
    inner_w = prs.slide_width - 2 * side_margin

    # ----- HEADER -----
    # Title block — centred serif, thin brass rules above and below.
    add_rect(slide, side_margin, Inches(0.85), inner_w, Inches(0.04), BRASS)
    tb, tf = add_text(slide, side_margin, Inches(1.05), inner_w, Inches(2.2),
                      anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    set_para(tf.paragraphs[0], TITLE, size=66, bold=True, color=INK,
             align=PP_ALIGN.CENTER, font=SERIF)
    # Byline italic
    tb, tf = add_text(slide, side_margin, Inches(3.3), inner_w, Inches(0.6),
                      anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    set_para(tf.paragraphs[0], VARIANT_AUTHOR_LINE, size=22, italic=True,
             color=MUTED, align=PP_ALIGN.CENTER, font=SERIF)
    add_rect(slide, side_margin, Inches(4.05), inner_w, Inches(0.04), BRASS)

    # ----- BODY: 3 columns -----
    col_top = Inches(4.5)
    col_bot = Inches(31.0)
    col_gap = Inches(0.7)
    col_w = (inner_w - 2 * col_gap) // 3
    col_x = [
        side_margin,
        side_margin + col_w + col_gap,
        side_margin + 2 * (col_w + col_gap),
    ]

    # ----- COLUMN 1: Abstract / Question / Background -----
    y = col_top
    _section_header(slide, col_x[0], y, col_w, "Abstract")
    y += Inches(0.95)
    abs_h = Inches(7.2)
    _panel(slide, col_x[0], y, col_w, abs_h)
    _body_prose(slide, col_x[0], y, col_w, abs_h, ABSTRACT, size=20, justify=True)
    y += abs_h + Inches(0.45)

    _section_header(slide, col_x[0], y, col_w, "Research question")
    y += Inches(0.95)
    q_h = Inches(3.0)
    _panel(slide, col_x[0], y, col_w, q_h)
    _body_prose(slide, col_x[0], y, col_w, q_h, QUESTION, size=20,
                justify=False, italic=True)
    y += q_h + Inches(0.45)

    _section_header(slide, col_x[0], y, col_w, "Background")
    y += Inches(0.95)
    bg_h = col_bot - y
    _panel(slide, col_x[0], y, col_w, bg_h)
    _bullets_serif(slide, col_x[0], y, col_w, bg_h - Inches(4.1),
                   BACKGROUND_BULLETS, size=16, leading=4)
    # Bottom of background panel: autocorrelation figure inset
    inset_y = y + bg_h - Inches(4.0)
    _figure_with_caption(slide, col_x[0] + Inches(0.2), inset_y,
                         col_w - Inches(0.4), Inches(3.2),
                         FIG_BACKGROUND,
                         "Figure 1. Label autocorrelation across lags "
                         "\u2014 motivates blocked CV.", cap_size=13)

    # ----- COLUMN 2: Methods prose + Headline figure + Take-home -----
    y = col_top
    _section_header(slide, col_x[1], y, col_w, "Methods")
    y += Inches(0.95)
    methods_h = Inches(8.0)
    _panel(slide, col_x[1], y, col_w, methods_h)
    methods_prose = " ".join(METHODS_BULLETS)
    _body_prose(slide, col_x[1], y, col_w, methods_h, methods_prose,
                size=17, justify=True)
    y += methods_h + Inches(0.45)

    _section_header(slide, col_x[1], y, col_w, "Headline figure")
    y += Inches(0.95)
    head_h = Inches(11.0)
    _panel(slide, col_x[1], y, col_w, head_h)
    _figure_with_caption(slide, col_x[1] + Inches(0.2), y + Inches(0.2),
                         col_w - Inches(0.4), head_h - Inches(1.0),
                         FIG_HEADLINE,
                         "Figure 2. KNN model selection under three "
                         "CV schemes.", cap_size=15)
    y += head_h + Inches(0.45)

    _section_header(slide, col_x[1], y, col_w, "Take-home")
    y += Inches(0.95)
    th_h = col_bot - y
    add_rect(slide, col_x[1], y, col_w, th_h, RGBColor(0xE9, 0xD8, 0xA8),
             line=BRASS, line_width_pt=0.6)
    _take_home_block(slide, col_x[1], y, col_w, th_h)

    # ----- COLUMN 3: Results / Conclusions / Supporting evidence -----
    y = col_top
    _section_header(slide, col_x[2], y, col_w, "Results")
    y += Inches(0.95)
    res_h = Inches(3.4)
    _results_table(slide, col_x[2], y, col_w, res_h)
    y += res_h + Inches(0.45)

    _section_header(slide, col_x[2], y, col_w, "Conclusions")
    y += Inches(0.95)
    conc_h = Inches(7.6)
    _panel(slide, col_x[2], y, col_w, conc_h)
    _bullets_serif(slide, col_x[2], y, col_w, conc_h,
                   CONCLUSION_BULLETS, size=15, leading=5)
    y += conc_h + Inches(0.45)

    _section_header(slide, col_x[2], y, col_w, "Supporting evidence")
    y += Inches(0.95)
    sup_h = col_bot - y
    _panel(slide, col_x[2], y, col_w, sup_h)
    _bullets_serif(slide, col_x[2], y, col_w, sup_h - Inches(4.1),
                   SUPPORTING_BULLETS, size=14, leading=3)
    inset_y = y + sup_h - Inches(4.0)
    _figure_with_caption(slide, col_x[2] + Inches(0.2), inset_y,
                         col_w - Inches(0.4), Inches(3.2),
                         FIG_SUPPORTING,
                         "Figure 3. Three-way CV comparison "
                         "\u2014 KNN shows the largest leakage gap.",
                         cap_size=13)

    # ----- REFERENCES FOOTER -----
    add_rect(slide, side_margin, Inches(31.5), inner_w, Inches(0.03), BRASS)
    _references_footer(slide, side_margin, Inches(31.65),
                       inner_w, Inches(3.7))

    # ----- COURSE ATTRIBUTION -----
    foot_y = Inches(35.45)
    tb, tf = add_text(slide, side_margin, foot_y, inner_w, Inches(0.45),
                      anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    set_para(tf.paragraphs[0],
             "COGS 109 \u2014 Modelling and Data Analysis \u2014 Spring 2026 "
             "\u2014 UC San Diego",
             size=13, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER, font=SERIF)

    set_core_properties(prs, subject="COGS 109 EEG poster \u2014 Academic Classic variant")
    prs.save(out_pptx)
    normalize_pptx_zip(out_pptx)
    render_preview_png(out_pptx, out_png)
    return out_pptx, out_png


if __name__ == "__main__":
    p, png = build()
    print(f"wrote {p}")
    print(f"wrote {png}")
