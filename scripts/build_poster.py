#!/usr/bin/env python3
"""Build the COGS 109 poster as a single-slide python-pptx file.

Produces ``poster/poster.pptx`` and ``poster/poster_preview.png``.

The design mirrors the NASA-style academic poster template Ivan reuses in
his Google Slides workspace: a dark navy banner header (title + author
strip + UCSD-style mark), three columns of content with blue section
header bars (Abstract / Introduction / Methods / Results / Conclusions /
References), a headline figure that anchors the centre column, and a
footer strip with course attribution.

Running this script twice produces a byte-identical ``poster.pptx`` so
the artifact is diffable in code review; the preview PNG is regenerated
each run.
"""

from __future__ import annotations

import datetime as _dt
import os
import shutil
import subprocess
import sys
import tempfile
import zipfile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
FIG_DIR = os.path.join(REPO_ROOT, "figures")
POSTER_DIR = os.path.join(REPO_ROOT, "poster")
POSTER_PPTX = os.path.join(POSTER_DIR, "poster.pptx")
POSTER_PNG = os.path.join(POSTER_DIR, "poster_preview.png")

# 48"x36" landscape academic poster
SLIDE_W_IN = 48.0
SLIDE_H_IN = 36.0

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
BAR_BLUE = RGBColor(0x1F, 0x3D, 0x7A)
ACCENT_GOLD = RGBColor(0xC6, 0x96, 0x14)
HIGHLIGHT_BG = RGBColor(0xF5, 0xEA, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x10, 0x10, 0x10)
LIGHT_BG = RGBColor(0xF7, 0xF7, 0xF7)
PANEL_BG = RGBColor(0xFA, 0xFA, 0xFC)

TITLE = "Classification of EEG Eye-State via KNN: A Cross-Validation Honesty Study"
AUTHORS = (
    "Ivan Del Rio  /  Anish Kondamadugula      "
    "COGS 109  •  Mukamel  •  Spring 2026      "
    "University of California, San Diego"
)
HEADER_KICKER = "UC San Diego  •  Department of Cognitive Science"

ABSTRACT = (
    "We pick classification as our data-analysis approach, K-nearest "
    "neighbours (KNN) as our model, and a k-sweep minimising 5-fold CV "
    "error as our model-selection procedure. We then run that same "
    "procedure under three different CV schemes — shuffled, naive blocked, "
    "and stratified blocked — on the UCI #264 EEG Eye State recording. "
    "All three schemes pick k = 1, but the picked accuracies span 47 "
    "percentage points: 97.3% ± 0.4% (shuffled, leaky), 50.0% ± 10.7% "
    "(naive blocked), and 77.8% ± 2.7% (stratified blocked, honest). The "
    "19.5 pp gap between the leaky and honest estimates is leakage "
    "attributable to scheme choice on this autocorrelated single-subject "
    "recording, not to model quality."
)

QUESTION = (
    "Given an autocorrelated single-subject EEG recording and a KNN "
    "classifier, how does the choice of cross-validation scheme change "
    "the result of an otherwise-identical model-selection procedure?"
)

BACKGROUND_BULLETS = [
    "UCI EEG Eye State dataset (Roesler, 2013) — id #264.",
    "14,980 samples × 14 Emotiv channels (AF3, F7, F3, FC5, T7, P7, O1, O2, "
    "P8, T8, FC6, F4, F8, AF4), sampled at ~128 Hz over ~117 s.",
    "Single subject, binary eyeDetection label, 55.12% open / 44.88% closed.",
    "Only 24 contiguous label runs in the whole recording.",
    "Channel autocorrelation at lag 1 ≈ 0.997 — adjacent samples are nearly "
    "identical, so any uniformly-random CV split puts near-duplicates into "
    "both train and test.",
]

METHODS_BULLETS = [
    "Approach: classification (binary eyeDetection target).",
    "Model: K-nearest neighbours with Euclidean distance in z-scored 14-D "
    "channel space. Chosen because KNN's per-sample similarity decision "
    "rule is the COGS 109-palette model most sensitive to lag-1 "
    "autocorrelation.",
    "Model selection: sweep k over {1, 3, 5, 7, 11, 15, 21, 31, 51, 75, 99, "
    "151, 201} (log-spaced); pick the k maximising mean 5-fold CV accuracy.",
    "CV scheme A — shuffled 5-fold (leaky baseline).",
    "CV scheme B — naive blocked 5-fold (5 contiguous time chunks).",
    "CV scheme C — stratified blocked 5-fold (100 short contiguous segments "
    "redistributed across folds to balance class proportion; honest).",
    "Preprocessing: drop 4 voltage outliers (>4 SD on any channel), "
    "chronological 80/20 split with a 64-sample seam gap, z-score on the "
    "training partition only.",
]

RESULTS_TABLE = [
    ["CV scheme", "Picked k", "Mean accuracy ± std"],
    ["Shuffled (leaky)", "1", "0.9728 ± 0.0037"],
    ["Naive blocked", "1", "0.5004 ± 0.1068"],
    ["Stratified blocked (honest)", "1", "0.7778 ± 0.0274"],
]

TAKEHOME_TOP = "Shuffled CV says we have a 97% classifier."
TAKEHOME_MID = "Honest stratified blocked CV says we have a 78% classifier."
TAKEHOME_BOT = "The difference is leakage, not signal."

CONCLUSION_BULLETS = [
    "(1) The same KNN-k-sweep model-selection procedure picks k = 1 under "
    "all three CV schemes — but the accuracy at that k varies by 47 "
    "percentage points (50.0% naive blocked → 97.3% shuffled).",
    "(2) The honest accuracy at k = 1 (stratified blocked CV) is 77.8% ± "
    "2.7% — well above the 55.12% majority-class baseline.",
    "(3) The 19.5 pp gap between the leaky and honest estimates is leakage "
    "attributable to scheme choice, not signal from a better model.",
    "(4) Alternative classifiers (LDA / PCA→LDA / PCR) confirm KNN is "
    "uniquely vulnerable — their leakage gaps are only 2–6 pp, as theory "
    "predicts for classifiers that do not depend on per-sample similarity.",
]

SUPPORTING_BULLETS = [
    "We also evaluated three alternative classifiers (LDA, PCA→LDA, "
    "PCR-as-classifier) as sanity checks, sweeping their hyperparameters "
    "under the same three CV schemes.",
    "LDA leakage gap (shuffled minus stratified blocked): +5.6 pp.",
    "PCA→LDA (n=3) leakage gap: +2.4 pp.",
    "PCR-as-classifier (n=2) leakage gap: +2.8 pp.",
    "KNN's +19.5 pp gap is the outlier — consistent with KNN exploiting "
    "per-sample similarity while LDA / PCA→LDA / PCR average over many "
    "samples before deciding. See figure 14 + tables/03 for full numbers.",
]

REFERENCES = [
    "[1] Roesler, O. (2013). EEG Eye State Data Set. UCI ML Repository #264. "
    "https://archive.ics.uci.edu/dataset/264/eeg+eye+state",
    "[2] Mukamel, R. (2026). COGS 109 — Modelling and Data Analysis "
    "Spring 2026 Study Guide. UC San Diego.",
    "[3] James, G., Witten, D., Hastie, T., Tibshirani, R. (2021). An "
    "Introduction to Statistical Learning, 2nd ed. Springer.",
    "[4] Bergmeir, C., Benítez, J. M. (2012). On the use of "
    "cross-validation for time series predictor evaluation. "
    "Information Sciences 191:192–213.",
]


def _add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def _set_para(p, text, *, size, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
              font="Calibri"):
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color


def _add_text(slide, x, y, w, h, *, fill=None, line=None, anchor=MSO_ANCHOR.TOP,
              margin=0.1):
    tb = slide.shapes.add_textbox(x, y, w, h)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    if line is not None:
        tb.line.color.rgb = line
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin / 2)
    tf.margin_bottom = Inches(margin / 2)
    # Clear default paragraph
    tf.paragraphs[0].text = ""
    return tb, tf


def _section_header(slide, x, y, w, label):
    """Blue bar + centred white section label."""
    _add_rect(slide, x, y, w, Inches(0.7), BAR_BLUE)
    tb, tf = _add_text(
        slide, x, y, w, Inches(0.7), anchor=MSO_ANCHOR.MIDDLE, margin=0.05
    )
    _set_para(tf.paragraphs[0], label, size=36, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER)


def _bullets(tf, items, *, size=20, color=BLACK, leading=4):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(leading)
        run = p.add_run()
        run.text = "•  " + item
        run.font.size = Pt(size)
        run.font.name = "Calibri"
        run.font.color.rgb = color


def _body_paragraph(tf, text, *, size=20, color=BLACK, justify=False, bold=False):
    p = tf.paragraphs[0] if not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = PP_ALIGN.JUSTIFY if justify else PP_ALIGN.LEFT
    p.space_after = Pt(6)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = "Calibri"
    run.font.bold = bold
    run.font.color.rgb = color


def _build_results_table(slide, x, y, w, h):
    rows, cols = len(RESULTS_TABLE), len(RESULTS_TABLE[0])
    tbl_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl = tbl_shape.table
    # Column widths: scheme name (wide) / picked k (narrow) / accuracy (wide).
    col_widths = [0.46, 0.18, 0.36]
    for i, frac in enumerate(col_widths):
        tbl.columns[i].width = int(w * frac)
    # Style every cell
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.06)
            cell.margin_bottom = Inches(0.06)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = BAR_BLUE
            elif RESULTS_TABLE[r][0].startswith("Stratified blocked"):
                cell.fill.fore_color.rgb = HIGHLIGHT_BG
            else:
                cell.fill.fore_color.rgb = PANEL_BG
            tf = cell.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = ""
            lines = RESULTS_TABLE[r][c].split("\n")
            for li, line in enumerate(lines):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                run = p.add_run()
                run.text = line
                run.font.size = Pt(18)
                run.font.name = "Calibri"
                # Bold the header row and the honest-row accuracy column.
                run.font.bold = (r == 0) or (
                    c == 2 and RESULTS_TABLE[r][0].startswith("Stratified blocked")
                )
                run.font.color.rgb = WHITE if r == 0 else BLACK
    return tbl_shape


def _add_figure_with_caption(slide, path, x, y, w, max_h, caption):
    """Add an image scaled to fit (w, max_h) plus a caption below."""
    from PIL import Image as _Image

    with _Image.open(path) as im:
        iw, ih = im.size
    aspect = ih / iw
    target_w = w
    target_h = int(target_w * aspect)
    if target_h > max_h:
        target_h = max_h
        target_w = int(max_h / aspect)
    offset_x = x + (w - target_w) // 2
    slide.shapes.add_picture(path, offset_x, y, width=target_w, height=target_h)
    cap_y = y + target_h + Inches(0.05)
    cap_h = Inches(0.45)
    tb, tf = _add_text(slide, x, cap_y, w, cap_h, anchor=MSO_ANCHOR.TOP,
                       margin=0.0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = caption
    run.font.size = Pt(16)
    run.font.italic = True
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
    return target_h + Inches(0.5).__int__() if hasattr(Inches(0.5), "__int__") else target_h


def build_poster():
    os.makedirs(POSTER_DIR, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # White background panel
    _add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE)

    # ----- HEADER -----
    header_h = Inches(5.0)
    _add_rect(slide, 0, 0, prs.slide_width, header_h, NAVY)
    # Accent stripe at the very bottom of the header
    _add_rect(slide, 0, header_h - Inches(0.18), prs.slide_width,
              Inches(0.18), ACCENT_GOLD)

    # Header text area (left of the UCSD mark)
    text_w = prs.slide_width - Inches(9.5)

    # Kicker line
    tb, tf = _add_text(slide, Inches(1.4), Inches(0.45),
                       text_w, Inches(0.6),
                       anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    _set_para(tf.paragraphs[0], HEADER_KICKER, size=28, bold=False,
              color=RGBColor(0xC9, 0xD2, 0xE6), align=PP_ALIGN.LEFT)

    # Title
    tb, tf = _add_text(slide, Inches(1.4), Inches(1.15),
                       text_w, Inches(2.6),
                       anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    _set_para(tf.paragraphs[0], TITLE, size=60, bold=True, color=WHITE,
              align=PP_ALIGN.LEFT)

    # Author strip (full width below the mark)
    tb, tf = _add_text(slide, Inches(1.4), Inches(3.95),
                       prs.slide_width - Inches(2.8), Inches(0.8),
                       anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    _set_para(tf.paragraphs[0], AUTHORS, size=26, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT)

    # UCSD-style mark on header right (text logo block, no external image).
    # Drawn AFTER the header so it appears on top.
    mark_w = Inches(7.0)
    mark_x = prs.slide_width - Inches(0.8) - mark_w
    mark_y = Inches(0.6)
    mark_h = Inches(3.4)
    _add_rect(slide, mark_x, mark_y, mark_w, mark_h, RGBColor(0x16, 0x2C, 0x52))
    tb, tf = _add_text(slide, mark_x, mark_y, mark_w, mark_h,
                       anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    _set_para(tf.paragraphs[0], "UC SAN DIEGO", size=44, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)
    p_sub = tf.add_paragraph()
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.space_before = Pt(8)
    r_sub = p_sub.add_run()
    r_sub.text = "COGS 109 — Spring 2026"
    r_sub.font.size = Pt(24)
    r_sub.font.bold = False
    r_sub.font.name = "Calibri"
    r_sub.font.color.rgb = ACCENT_GOLD

    # ----- Layout: 3 columns -----
    margin_x = Inches(1.0)
    gap = Inches(0.6)
    col_w = (prs.slide_width - 2 * margin_x - 2 * gap) // 3
    col_top = header_h + Inches(0.6)
    col_bottom = prs.slide_height - Inches(1.4)
    col_h = col_bottom - col_top
    col_x = [
        margin_x,
        margin_x + col_w + gap,
        margin_x + 2 * (col_w + gap),
    ]

    # ----- LEFT COLUMN -----
    y = col_top

    _section_header(slide, col_x[0], y, col_w, "ABSTRACT")
    y += Inches(0.85)
    abstract_h = Inches(6.6)
    tb, tf = _add_text(slide, col_x[0], y, col_w, abstract_h,
                       fill=PANEL_BG, anchor=MSO_ANCHOR.TOP, margin=0.25)
    _body_paragraph(tf, ABSTRACT, size=20, justify=True)
    y += abstract_h + Inches(0.35)

    _section_header(slide, col_x[0], y, col_w, "QUESTION")
    y += Inches(0.85)
    q_h = Inches(2.6)
    tb, tf = _add_text(slide, col_x[0], y, col_w, q_h,
                       fill=HIGHLIGHT_BG, anchor=MSO_ANCHOR.MIDDLE, margin=0.25)
    _body_paragraph(tf, QUESTION, size=22, justify=True, bold=True)
    y += q_h + Inches(0.35)

    _section_header(slide, col_x[0], y, col_w, "BACKGROUND")
    y += Inches(0.85)
    bg_h = col_bottom - y
    tb, tf = _add_text(slide, col_x[0], y, col_w, bg_h, fill=PANEL_BG,
                       anchor=MSO_ANCHOR.TOP, margin=0.25)
    _bullets(tf, BACKGROUND_BULLETS, size=18)
    # Embed the autocorrelation figure inside Background panel — bottom half
    auto_y = y + Inches(4.0)
    _add_figure_with_caption(
        slide,
        os.path.join(FIG_DIR, "08_label_autocorrelation.png"),
        col_x[0] + Inches(0.2), auto_y, col_w - Inches(0.4),
        bg_h - Inches(4.1),
        "Fig. 1.  Label autocorrelation across lags — motivates blocked CV.",
    )

    # ----- MIDDLE COLUMN -----
    y = col_top
    _section_header(slide, col_x[1], y, col_w, "METHODS")
    y += Inches(0.85)
    methods_h = Inches(6.0)
    tb, tf = _add_text(slide, col_x[1], y, col_w, methods_h, fill=PANEL_BG,
                       anchor=MSO_ANCHOR.TOP, margin=0.25)
    _bullets(tf, METHODS_BULLETS, size=18)
    y += methods_h + Inches(0.35)

    _section_header(slide, col_x[1], y, col_w, "TAKE-HOME")
    y += Inches(0.85)
    take_h = Inches(4.6)
    tb, tf = _add_text(slide, col_x[1], y, col_w, take_h, fill=HIGHLIGHT_BG,
                       anchor=MSO_ANCHOR.MIDDLE, margin=0.3)
    _set_para(tf.paragraphs[0], TAKEHOME_TOP, size=28, bold=True, color=BLACK,
              align=PP_ALIGN.CENTER)
    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(8)
    r1 = p1.add_run()
    r1.text = TAKEHOME_MID
    r1.font.size = Pt(28)
    r1.font.bold = True
    r1.font.name = "Calibri"
    r1.font.color.rgb = RGBColor(0xB0, 0x30, 0x30)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    r2 = p2.add_run()
    r2.text = TAKEHOME_BOT
    r2.font.size = Pt(24)
    r2.font.italic = True
    r2.font.name = "Calibri"
    r2.font.color.rgb = BLACK
    y += take_h + Inches(0.35)

    _section_header(slide, col_x[1], y, col_w, "HEADLINE FIGURE")
    y += Inches(0.85)
    head_fig_h = col_bottom - y
    _add_figure_with_caption(
        slide,
        os.path.join(FIG_DIR, "11_knn_k_sweep.png"),
        col_x[1], y, col_w, head_fig_h - Inches(0.6),
        "Fig. 2.  KNN mean accuracy ± std-dev vs k (log scale) under three "
        "CV schemes — all three schemes pick k = 1 but report accuracies "
        "spanning 47 percentage points.",
    )

    # ----- RIGHT COLUMN -----
    y = col_top
    _section_header(slide, col_x[2], y, col_w, "RESULTS")
    y += Inches(0.85)
    tbl_h = Inches(3.4)
    _build_results_table(slide, col_x[2], y, col_w, tbl_h)
    y += tbl_h + Inches(0.35)

    # Supporting evidence panel: alternative classifiers (LDA / PCA->LDA /
    # PCR) plus the secondary figure 14 (4-model 3-scheme bar chart) and
    # figure 15 (holdout confusion matrices, interpret with care).
    _section_header(slide, col_x[2], y, col_w, "SUPPORTING EVIDENCE")
    y += Inches(0.85)
    sup_text_h = Inches(3.6)
    tb, tf = _add_text(slide, col_x[2], y, col_w, sup_text_h, fill=PANEL_BG,
                       anchor=MSO_ANCHOR.TOP, margin=0.25)
    _bullets(tf, SUPPORTING_BULLETS, size=16)
    y += sup_text_h + Inches(0.2)

    sup_fig_h = Inches(3.0)
    _add_figure_with_caption(
        slide,
        os.path.join(FIG_DIR, "14_cv_comparison_three_way.png"),
        col_x[2], y, col_w, sup_fig_h - Inches(0.5),
        "Fig. 3.  Supporting figure — three-way CV comparison across LDA "
        "/ KNN / PCA→LDA / PCR. KNN bar group has the largest leakage gap.",
    )
    y += sup_fig_h + Inches(0.2)

    _section_header(slide, col_x[2], y, col_w, "CONCLUSIONS")
    y += Inches(0.85)
    conc_h = Inches(4.4)
    tb, tf = _add_text(slide, col_x[2], y, col_w, conc_h, fill=PANEL_BG,
                       anchor=MSO_ANCHOR.TOP, margin=0.25)
    _bullets(tf, CONCLUSION_BULLETS, size=15)
    y += conc_h + Inches(0.25)

    _section_header(slide, col_x[2], y, col_w, "REFERENCES")
    y += Inches(0.85)
    ref_h = col_bottom - y
    tb, tf = _add_text(slide, col_x[2], y, col_w, ref_h, fill=PANEL_BG,
                       anchor=MSO_ANCHOR.TOP, margin=0.25)
    for i, ref in enumerate(REFERENCES):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = ref
        run.font.size = Pt(14)
        run.font.name = "Calibri"
        run.font.color.rgb = BLACK

    # ----- FOOTER -----
    footer_y = prs.slide_height - Inches(1.0)
    _add_rect(slide, 0, footer_y, prs.slide_width, Inches(1.0), NAVY)
    tb, tf = _add_text(slide, Inches(1.4), footer_y, prs.slide_width - Inches(2.8),
                       Inches(1.0), anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    _set_para(tf.paragraphs[0],
              "Source code & full report:  "
              "github.com/BariBariGood/cogs109-eeg-eye-state-classifier",
              size=20, color=RGBColor(0xD8, 0xDE, 0xE8), align=PP_ALIGN.LEFT)

    # Right-side credit on the footer
    tb, tf = _add_text(slide, prs.slide_width - Inches(15),
                       footer_y, Inches(13.6),
                       Inches(1.0), anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    _set_para(tf.paragraphs[0],
              "COGS 109 — Modelling and Data Analysis — Spring 2026",
              size=20, color=RGBColor(0xD8, 0xDE, 0xE8), align=PP_ALIGN.RIGHT)

    # Pin a deterministic "created" property so byte output is stable.
    cp = prs.core_properties
    cp.author = "BariBariGood"
    cp.title = TITLE
    cp.subject = "COGS 109 Spring 2026 — EEG Eye-State Classifier"
    fixed = _dt.datetime(2026, 5, 28, 0, 0, 0)
    cp.created = fixed
    cp.modified = fixed
    cp.last_modified_by = "BariBariGood"

    prs.save(POSTER_PPTX)
    _normalize_pptx_zip(POSTER_PPTX)
    return POSTER_PPTX


def _normalize_pptx_zip(path):
    """Re-pack ``path`` with deterministic ZIP entry order + fixed timestamps.

    python-pptx delegates to ``zipfile.ZipFile`` which captures the wall-clock
    time of each write. We re-open the archive, sort entries by name, and
    rewrite every entry with the same (epoch-style) timestamp so two runs of
    ``build_poster.py`` produce byte-identical .pptx files.
    """
    fixed_date = (2026, 5, 28, 0, 0, 0)
    with zipfile.ZipFile(path, "r") as src:
        infos = sorted(src.infolist(), key=lambda zi: zi.filename)
        payloads = [(zi.filename, src.read(zi.filename)) for zi in infos]

    tmp_path = path + ".tmp"
    with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as dst:
        for name, data in payloads:
            zi = zipfile.ZipInfo(filename=name, date_time=fixed_date)
            zi.compress_type = zipfile.ZIP_DEFLATED
            zi.external_attr = 0o644 << 16
            dst.writestr(zi, data)
    os.replace(tmp_path, path)


def render_preview_png():
    """Render a PNG preview of ``poster.pptx``.

    The preferred path is LibreOffice's headless ``impress_png_Export`` filter,
    which renders the .pptx as-is (deterministic byte output across runs on
    the same LO version). If ``soffice`` is unavailable, we fall back to a
    minimal PIL replica that at least communicates the layout.
    """
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice is not None:
        with tempfile.TemporaryDirectory(prefix="poster_render_") as td:
            try:
                subprocess.run(
                    [soffice, "--headless", "--convert-to", "png",
                     "--outdir", td, POSTER_PPTX],
                    check=True, capture_output=True, timeout=120,
                )
                rendered = os.path.join(td, "poster.png")
                if os.path.exists(rendered):
                    shutil.copyfile(rendered, POSTER_PNG)
                    return POSTER_PNG
            except (subprocess.CalledProcessError, subprocess.TimeoutExpired):
                pass
    # Fallback: PIL stub preview
    return _pil_fallback_preview()


def _pil_fallback_preview():
    """Minimal PIL preview when LibreOffice is not available.

    Renders the header banner + a centred placeholder telling the reader
    to open the .pptx directly. Not pretty, but never crashes the build.
    """
    scale = 30  # pixels per inch
    W = int(SLIDE_W_IN * scale)
    H = int(SLIDE_H_IN * scale)
    img = Image.new("RGB", (W, H), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    def _font(size):
        for candidate in (
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        ):
            if os.path.exists(candidate):
                return ImageFont.truetype(candidate, size)
        return ImageFont.load_default()

    def _font_bold(size):
        candidate = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
        if os.path.exists(candidate):
            return ImageFont.truetype(candidate, size)
        return ImageFont.load_default()

    def px(v):
        return int(v * scale)

    # Header
    header_h = px(5.0)
    draw.rectangle([0, 0, W, header_h], fill=(11, 31, 58))
    draw.rectangle([0, header_h - px(0.18), W, header_h], fill=(198, 150, 20))
    draw.text((px(1.4), px(0.55)), HEADER_KICKER,
              fill=(201, 210, 230), font=_font(int(0.6 * scale)))
    # Title — wrap manually
    draw.text((px(1.4), px(1.4)),
              "Classification of EEG Eye-State via KNN:",
              fill=(255, 255, 255), font=_font_bold(int(1.3 * scale)))
    draw.text((px(1.4), px(2.5)),
              "A Cross-Validation Honesty Study",
              fill=(255, 255, 255), font=_font_bold(int(1.3 * scale)))
    draw.text((px(1.4), px(3.95)),
              AUTHORS,
              fill=(240, 240, 240), font=_font(int(0.55 * scale)))
    # UCSD mark on the right of the header
    mark_x = W - px(7.5)
    mark_y = px(0.6)
    draw.rectangle([mark_x, mark_y, mark_x + px(6.4), mark_y + px(3.6)],
                   fill=(22, 44, 82))
    mark_title = "UC SAN DIEGO"
    mark_font = _font_bold(int(1.0 * scale))
    mtw = draw.textlength(mark_title, font=mark_font)
    draw.text((mark_x + (px(6.4) - mtw) // 2, mark_y + px(0.9)),
              mark_title, fill=(255, 255, 255), font=mark_font)
    mark_sub = "COGS 109 — Spring 2026"
    msub_font = _font(int(0.55 * scale))
    mstw = draw.textlength(mark_sub, font=msub_font)
    draw.text((mark_x + (px(6.4) - mstw) // 2, mark_y + px(2.3)),
              mark_sub, fill=(198, 150, 20), font=msub_font)

    # 3 columns
    margin_x = px(1.0)
    gap = px(0.6)
    col_w = (W - 2 * margin_x - 2 * gap) // 3
    col_top = header_h + px(0.6)
    col_bottom_px = H - px(1.4)
    col_x = [margin_x, margin_x + col_w + gap, margin_x + 2 * (col_w + gap)]

    def _section(col, y, label):
        x0 = col_x[col]
        draw.rectangle([x0, y, x0 + col_w, y + px(0.7)], fill=(31, 61, 122))
        text = label
        tw = draw.textlength(text, font=_font_bold(int(0.7 * scale)))
        tx = x0 + (col_w - tw) // 2
        draw.text((tx, y + px(0.08)), text, fill=(255, 255, 255),
                  font=_font_bold(int(0.7 * scale)))
        return y + px(0.85)

    def _panel(col, y, h, fill=(250, 250, 252)):
        x0 = col_x[col]
        draw.rectangle([x0, y, x0 + col_w, y + h], fill=fill)
        return x0, y, x0 + col_w, y + h

    def _wrap_text(text, font, max_w):
        words = text.split()
        lines, cur = [], ""
        for word in words:
            trial = (cur + " " + word).strip()
            if draw.textlength(trial, font=font) <= max_w:
                cur = trial
            else:
                if cur:
                    lines.append(cur)
                cur = word
        if cur:
            lines.append(cur)
        return lines

    # ---- LEFT COLUMN ----
    y = col_top
    y = _section(0, y, "ABSTRACT")
    abstract_h = px(7.0)
    _panel(0, y, abstract_h)
    body_font = _font(int(0.5 * scale))
    lines = _wrap_text(ABSTRACT, body_font, col_w - px(0.5))
    cy = y + px(0.25)
    for ln in lines:
        draw.text((col_x[0] + px(0.25), cy), ln, fill=(20, 20, 20),
                  font=body_font)
        cy += int(0.62 * scale)
    y += abstract_h + px(0.4)

    y = _section(0, y, "QUESTION")
    q_h = px(2.6)
    _panel(0, y, q_h, fill=(245, 234, 200))
    qfont = _font_bold(int(0.55 * scale))
    lines = _wrap_text(QUESTION, qfont, col_w - px(0.5))
    cy = y + px(0.25)
    for ln in lines:
        draw.text((col_x[0] + px(0.25), cy), ln, fill=(20, 20, 20),
                  font=qfont)
        cy += int(0.7 * scale)
    y += q_h + px(0.4)

    y = _section(0, y, "BACKGROUND")
    bg_h_left = H - px(1.4) - y  # to footer
    _panel(0, y, bg_h_left)
    bf = _font(int(0.46 * scale))
    cy = y + px(0.25)
    for item in BACKGROUND_BULLETS:
        wrapped = _wrap_text("• " + item, bf, col_w - px(0.5))
        for ln in wrapped:
            draw.text((col_x[0] + px(0.25), cy), ln, fill=(20, 20, 20),
                      font=bf)
            cy += int(0.55 * scale)
        cy += int(0.1 * scale)

    # Inline autocorrelation figure on the bottom half of the background panel
    auto_path = os.path.join(FIG_DIR, "08_label_autocorrelation.png")
    if os.path.exists(auto_path):
        auto = Image.open(auto_path).convert("RGB")
        target_w = col_w - px(0.5)
        target_h = int(auto.height * target_w / auto.width)
        max_fig_h = bg_h_left - (cy - y) - px(0.6)
        if target_h > max_fig_h and max_fig_h > 0:
            target_h = max_fig_h
            target_w = int(auto.width * target_h / auto.height)
        if target_w > 0 and target_h > 0:
            auto = auto.resize((target_w, target_h), Image.LANCZOS)
            ax = col_x[0] + (col_w - target_w) // 2
            img.paste(auto, (ax, cy + px(0.2)))
            cap_y = cy + px(0.2) + target_h + px(0.1)
            cap_font = _font(int(0.4 * scale))
            cap = "Fig. 1.  Label autocorrelation across lags."
            draw.text((col_x[0] + px(0.25), cap_y), cap,
                      fill=(60, 60, 60), font=cap_font)

    # ---- MIDDLE COLUMN ----
    y = col_top
    y = _section(1, y, "METHODS")
    m_h = px(7.0)
    _panel(1, y, m_h)
    cy = y + px(0.25)
    for item in METHODS_BULLETS:
        wrapped = _wrap_text("• " + item, bf, col_w - px(0.5))
        for ln in wrapped:
            draw.text((col_x[1] + px(0.25), cy), ln, fill=(20, 20, 20),
                      font=bf)
            cy += int(0.55 * scale)
        cy += int(0.1 * scale)
    y += m_h + px(0.4)

    y = _section(1, y, "TAKE-HOME")
    th_h = px(3.4)
    _panel(1, y, th_h, fill=(245, 234, 200))
    big = _font_bold(int(0.85 * scale))
    bigred = _font_bold(int(0.85 * scale))
    italic = _font(int(0.75 * scale))
    lines_top = _wrap_text(TAKEHOME_TOP, big, col_w - px(0.5))
    lines_mid = _wrap_text(TAKEHOME_MID, big, col_w - px(0.5))
    lines_bot = _wrap_text(TAKEHOME_BOT, italic, col_w - px(0.5))
    cy = y + px(0.25)
    for ln in lines_top:
        tw = draw.textlength(ln, font=big)
        draw.text((col_x[1] + (col_w - tw) // 2, cy), ln,
                  fill=(20, 20, 20), font=big)
        cy += int(1.0 * scale)
    cy += int(0.1 * scale)
    for ln in lines_mid:
        tw = draw.textlength(ln, font=bigred)
        draw.text((col_x[1] + (col_w - tw) // 2, cy), ln,
                  fill=(176, 48, 48), font=bigred)
        cy += int(1.0 * scale)
    for ln in lines_bot:
        tw = draw.textlength(ln, font=italic)
        draw.text((col_x[1] + (col_w - tw) // 2, cy), ln,
                  fill=(20, 20, 20), font=italic)
        cy += int(0.85 * scale)
    y += th_h + px(0.4)

    y = _section(1, y, "HEADLINE FIGURE")
    head_fig_path = os.path.join(FIG_DIR, "11_knn_k_sweep.png")
    if os.path.exists(head_fig_path):
        fig = Image.open(head_fig_path).convert("RGB")
        target_w = col_w
        target_h = int(fig.height * target_w / fig.width)
        avail_h = (H - px(1.4)) - y - px(0.6)
        if target_h > avail_h and avail_h > 0:
            target_h = avail_h
            target_w = int(fig.width * target_h / fig.height)
        if target_w > 0 and target_h > 0:
            fig = fig.resize((target_w, target_h), Image.LANCZOS)
            fx = col_x[1] + (col_w - target_w) // 2
            img.paste(fig, (fx, y))
            cap_font = _font(int(0.4 * scale))
            draw.text((col_x[1] + px(0.1), y + target_h + px(0.1)),
                      "Fig. 2. KNN k-sweep under three CV schemes.",
                      fill=(60, 60, 60), font=cap_font)

    # ---- RIGHT COLUMN ----
    y = col_top
    y = _section(2, y, "RESULTS")
    # Render the results table manually
    rows = len(RESULTS_TABLE)
    cols = len(RESULTS_TABLE[0])
    tbl_h = px(3.4)
    row_h = tbl_h // rows
    col_fracs = [0.46, 0.18, 0.36]
    col_widths = [int(col_w * f) for f in col_fracs]
    table_font = _font_bold(int(0.45 * scale))
    table_font_b = _font(int(0.42 * scale))
    table_font_hl = _font_bold(int(0.45 * scale))
    for r in range(rows):
        cx = col_x[2]
        for c in range(cols):
            x0 = cx
            x1 = cx + col_widths[c]
            y0 = y + r * row_h
            y1 = y + (r + 1) * row_h
            is_honest = RESULTS_TABLE[r][0].startswith("Stratified blocked")
            if r == 0:
                fill = (31, 61, 122)
                fg = (255, 255, 255)
            elif is_honest:
                fill = (245, 234, 200)
                fg = (20, 20, 20)
            else:
                fill = (250, 250, 252)
                fg = (20, 20, 20)
            draw.rectangle([x0, y0, x1, y1], fill=fill,
                           outline=(200, 200, 200))
            text = RESULTS_TABLE[r][c]
            if r == 0:
                f = table_font
            elif c == 2 and is_honest:
                f = table_font_hl
            else:
                f = table_font_b
            text_lines = text.split("\n")
            total_h = int(0.55 * scale) * len(text_lines)
            line_y = y0 + (row_h - total_h) // 2
            for ln in text_lines:
                tw = draw.textlength(ln, font=f)
                tx = x0 + (col_widths[c] - tw) // 2 if c > 0 else x0 + px(0.15)
                draw.text((tx, line_y), ln, fill=fg, font=f)
                line_y += int(0.55 * scale)
            cx = x1
    y += tbl_h + px(0.35)

    y = _section(2, y, "SUPPORTING EVIDENCE")
    sup_h = px(3.6)
    _panel(2, y, sup_h)
    sf = _font(int(0.4 * scale))
    cy = y + px(0.2)
    for item in SUPPORTING_BULLETS:
        wrapped = _wrap_text("• " + item, sf, col_w - px(0.5))
        for ln in wrapped:
            draw.text((col_x[2] + px(0.25), cy), ln, fill=(20, 20, 20),
                      font=sf)
            cy += int(0.46 * scale)
        cy += int(0.08 * scale)
    y += sup_h + px(0.25)

    aux_path = os.path.join(FIG_DIR, "14_cv_comparison_three_way.png")
    aux_h = px(3.0)
    if os.path.exists(aux_path):
        aux = Image.open(aux_path).convert("RGB")
        target_w = col_w
        target_h = int(aux.height * target_w / aux.width)
        if target_h > aux_h - px(0.5):
            target_h = aux_h - px(0.5)
            target_w = int(aux.width * target_h / aux.height)
        aux = aux.resize((target_w, target_h), Image.LANCZOS)
        fx = col_x[2] + (col_w - target_w) // 2
        img.paste(aux, (fx, y))
        cap_font = _font(int(0.36 * scale))
        draw.text((col_x[2] + px(0.1), y + target_h + px(0.05)),
                  "Fig. 3. Supporting — three-way CV across 4 classifiers.",
                  fill=(60, 60, 60), font=cap_font)
    y += aux_h + px(0.2)

    y = _section(2, y, "CONCLUSIONS")
    conc_h = px(5.0)
    _panel(2, y, conc_h)
    cf = _font(int(0.38 * scale))
    cy = y + px(0.2)
    for item in CONCLUSION_BULLETS:
        wrapped = _wrap_text("• " + item, cf, col_w - px(0.5))
        for ln in wrapped:
            draw.text((col_x[2] + px(0.25), cy), ln, fill=(20, 20, 20),
                      font=cf)
            cy += int(0.46 * scale)
        cy += int(0.08 * scale)
    y += conc_h + px(0.3)

    y = _section(2, y, "REFERENCES")
    ref_h = (H - px(1.4)) - y
    _panel(2, y, ref_h)
    rf = _font(int(0.36 * scale))
    cy = y + px(0.2)
    for ref in REFERENCES:
        wrapped = _wrap_text(ref, rf, col_w - px(0.5))
        for ln in wrapped:
            draw.text((col_x[2] + px(0.25), cy), ln, fill=(20, 20, 20),
                      font=rf)
            cy += int(0.42 * scale)
        cy += int(0.08 * scale)

    # Footer
    footer_y = H - px(1.0)
    draw.rectangle([0, footer_y, W, H], fill=(11, 31, 58))
    foot_font = _font(int(0.5 * scale))
    draw.text((px(1.4), footer_y + px(0.3)),
              "Source code & full report: "
              "github.com/BariBariGood/cogs109-eeg-eye-state-classifier",
              fill=(216, 222, 232), font=foot_font)
    credit = "COGS 109 — Modelling and Data Analysis — Spring 2026"
    tw = draw.textlength(credit, font=foot_font)
    draw.text((W - px(1.4) - tw, footer_y + px(0.3)), credit,
              fill=(216, 222, 232), font=foot_font)

    img.save(POSTER_PNG, "PNG", optimize=True)
    return POSTER_PNG


def main():
    pptx_path = build_poster()
    png_path = render_preview_png()
    print(f"wrote {pptx_path}")
    print(f"wrote {png_path}")


if __name__ == "__main__":
    main()
